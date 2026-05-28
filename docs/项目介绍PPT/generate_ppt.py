#!/usr/bin/env python3
"""Generate BloodLine-Analysis project introduction PPT."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x7E)
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_section_slide(prs, title):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x7E)
    shape.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullets, accent_color=RGBColor(0x1A, 0x23, 0x7E)):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"●  {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(14)
        p.level = 0

    return slide


def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x7E)
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.3), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    # Left column items
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.3), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"●  {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(10)

    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    # Right column items
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.3), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"●  {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(10)

    return slide


def add_architecture_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x7E)
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "系统架构"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    # Layer boxes
    layers = [
        ("前端展示层", "React + Vite | 表级血缘图谱 | 调度查看 | 模糊搜索", RGBColor(0xE3, 0xF2, 0xFD), Inches(0.5)),
        ("API 服务层", "FastAPI | 血缘扫描引擎 | 血缘查询服务 | Excel导出 | MySQL同步", RGBColor(0xE8, 0xF5, 0xE9), Inches(2.0)),
        ("数据解析层", "Kettle Repo解析 | Java代码SQL提取 | FineReport SQL解析 | MySQL元数据采集", RGBColor(0xFF, 0xF3, 0xE0), Inches(3.5)),
        ("数据源层", "Kettle作业仓库(.repo) | Java源码 | MySQL视图定义 | FineReport报表库", RGBColor(0xF3, 0xE5, 0xF5), Inches(5.0)),
    ]

    for label, desc, color, top in layers:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top, Inches(9), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = RGBColor(0x1A, 0x23, 0x7E)
        box.line.width = Pt(1.5)

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{label}\n{desc}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.alignment = PP_ALIGN.CENTER

    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. 封面
    add_title_slide(prs, "BloodLine-Analysis", "数据血缘分析与可视化平台")

    # 2. 目录
    add_content_slide(prs, "目录", [
        "项目背景与挑战",
        "系统架构设计",
        "核心功能介绍",
        "技术实现亮点",
        "应用价值与成果",
        "未来发展规划",
    ])

    # 3. 项目背景
    add_content_slide(prs, "项目背景与挑战", [
        "数据链路复杂：ETL作业（Kettle）、Java代码、FineReport报表多方参与数据处理",
        "血缘关系黑盒化：表与表之间的数据流向难以追踪，影响影响分析和故障排查",
        "元数据分散：MySQL 6个业务库（DM/DP/DPINC/FRMS/HQ/SG）的表结构和视图定义缺乏统一视图",
        "缺乏可视化工具：开发人员无法直观查看数据从哪来、到哪去、经过了哪些加工",
        "数据治理需求：需要支持血缘同步到 MySQL DM.t_relationship 供下游系统消费",
    ])

    # 4. 系统架构
    add_architecture_slide(prs)

    # 5. 核心功能 - 血缘采集
    add_two_column_slide(prs,
        "核心功能 — 多源血缘采集",
        "结构化数据源",
        [
            "Kettle ETL 作业：解析 .repo XML 文件，提取 Transformation Step 和 Job Entry 的 SQL",
            "MySQL 元数据：扫描 6 个业务库，采集表结构、视图定义（VIEW）",
            "FineReport 报表：读取 frms.comm_finereport_record_details，解析 data_sql 中的来源表",
        ],
        "代码数据源",
        [
            "Java 源码分析：静态解析 Java 文件中的 SQL 语句，提取读写表关系",
            "SQL 智能提取：支持 SELECT/INSERT/UPDATE/DELETE 多类型语句的表名识别",
            "字段级血缘：解析 SQL 列映射关系，构建字段级 READS/WRITES 边",
        ]
    )

    # 6. 核心功能 - 血缘构建
    add_content_slide(prs, "核心功能 — 血缘关系构建", [
        "Actor-Scope 配对机制：同一 Transformation 内跨 Step 配对 READS/WRITES",
        "Job Entry 级隔离：同一 Job 内不同 Entry 的 READS/WRITES 互不干扰",
        "FLOWS_TO 推导：基于共享 Actor 推导表级数据流（Table Flow）",
        "视图血缘：自动解析 MySQL 视图定义，建立视图 → 源表的 FLOWS_TO 关系",
        "报表血缘：FineReport 数据集与来源表建立 FLOWS_TO，标记 source=finereport",
        "图数据模型：Node（表/作业/模块）+ Edge（READS/WRITES/CALLS/FLOWS_TO）",
    ])

    # 7. 核心功能 - 前端展示
    add_two_column_slide(prs,
        "核心功能 — 前端可视化",
        "血缘图谱",
        [
            "表详情页：展示单张表的上游来源和下游去向",
            "自环检测：自动发现表指向自身的异常血缘",
            "闭环检测：识别多表循环依赖链路",
        ],
        "调度查看",
        [
            "血缘列表：平铺展示所有表级血缘关系",
            "模糊搜索：支持按来源表名、目标表名模糊过滤",
            "Actor 归因：自动标注血缘的加工方（Kettle Job / Java 模块 / FineReport）",
        ]
    )

    # 8. 核心功能 - 血缘同步
    add_content_slide(prs, "核心功能 — 血缘同步与导出", [
        "MySQL 血缘同步：将 SQLite FLOWS_TO 边批量写入 DM.t_relationship",
        "Excel 导出：支持全量血缘关系导出为 Excel 报表",
        "扫描历史：记录每次扫描的输入参数、状态、耗时，支持回溯",
        "失败追踪：扫描过程中记录解析失败（SQL 语法错误、文件缺失等）",
    ])

    # 9. 技术亮点
    add_two_column_slide(prs,
        "技术实现亮点",
        "解析引擎",
        [
            "多层 Actor Scope：transformation 级 collapse vs job entry 级隔离",
            "Schema 感知：统一规范化表名（schema.table），处理大小写差异",
            "别名解析：基于 MySQL 元数据建立表名别名映射，提高识别准确率",
        ],
        "工程实践",
        [
            "FastAPI + SQLAlchemy：异步 ORM，支持事务性批量写入",
            "Graph State 重置：每次全量扫描前清理旧图，保证数据一致性",
            "模块化设计：采集/构建/查询/导出/同步各司其职，易于扩展",
        ]
    )

    # 10. 应用价值
    add_content_slide(prs, "应用价值", [
        "影响分析：数据表变更时，快速定位所有下游依赖，降低变更风险",
        "故障排查：数据异常时，追溯数据来源链路，定位问题根因",
        "数据治理：建立统一的数据资产地图，辅助数据分类分级",
        "合规审计：提供完整的数据流转记录，满足数据安全审计要求",
        "知识沉淀：将分散在 ETL 脚本和代码中的隐性知识显性化、可视化",
    ])

    # 11. 未来规划
    add_content_slide(prs, "未来发展规划", [
        "增量扫描：支持按时间戳或变更集进行增量血缘更新，提升扫描效率",
        "字段级血缘可视化：在前端展示列级别的数据流向",
        "血缘质量评分：基于覆盖率、准确率等指标评估血缘数据质量",
        "告警订阅：数据链路断裂或异常时自动推送告警通知",
        "多数据源扩展：支持 Hive、ClickHouse、Oracle 等更多数据库的血缘采集",
    ])

    # 12. 结束页
    add_title_slide(prs, "谢谢", "BloodLine-Analysis 数据血缘平台")

    output_path = "/Users/fengxiaomao/BloodLine-Analysis/docs/项目介绍PPT/BloodLine-Analysis项目介绍.pptx"
    prs.save(output_path)
    print(f"PPT generated: {output_path}")


if __name__ == "__main__":
    main()
